#!/usr/bin/env python3
"""Export an architecture-diagram SVG (or HTML) to a native, editable PPTX.

Wraps the vendored svg_to_pptx converter (from ppt-master): each SVG
element becomes a native DrawingML shape, so rounded corners, dashed
borders and arrowheads survive as editable PowerPoint shapes.
(PowerPoint's own SVG-import -> convert-to-shape path drops them.)

Requires python-pptx:  pip install -r requirements.txt
"""

from __future__ import annotations

import argparse
import re
import shutil
import sys
import tempfile
import zipfile
from pathlib import Path

# Make the sibling svg_to_pptx package importable when run as a script.
SCRIPTS_DIR = Path(__file__).resolve().parent
if str(SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPTS_DIR))

from svg_to_pptx import convert_svg_to_slide_shapes  # noqa: E402

EMU_PER_PX = 9525
SVG_TAG_RE = re.compile(r"<svg\b.*?</svg>", re.IGNORECASE | re.DOTALL)


def _viewbox_size(svg_text: str) -> tuple[float, float]:
    """Return (width, height) in SVG px from the root viewBox."""
    m = re.search(r"<svg\b[^>]*\bviewBox\s*=\s*[\"']([^\"']+)[\"']", svg_text, re.IGNORECASE)
    if not m:
        raise ValueError("SVG has no viewBox; cannot size the slide")
    parts = m.group(1).replace(",", " ").split()
    if len(parts) != 4:
        raise ValueError(f"Unexpected viewBox: {m.group(1)}")
    return float(parts[2]), float(parts[3])


def _preprocess_svg(svg_text: str, w: float, h: float) -> str:
    """Normalize SVG for the converter: resolve 100% width/height to viewBox px.

    The converter parses width/height as floats; '100%' (used by the
    full-canvas background rects) would crash it. The grid <pattern> fill
    is left as-is (renders empty in PPTX - decorative, acceptable loss).
    """
    svg_text = re.sub(r'\bwidth\s*=\s*"100%"', f'width="{w:g}"', svg_text)
    svg_text = re.sub(r'\bheight\s*=\s*"100%"', f'height="{h:g}"', svg_text)
    return svg_text


def _extract_svg_from_html(html_text: str) -> str:
    m = SVG_TAG_RE.search(html_text)
    if not m:
        raise ValueError("No <svg> block found in HTML")
    return m.group(0)


def convert_svg_to_pptx(svg_path: Path, out_path: Path, verbose: bool = False) -> Path:
    """Convert one SVG file to a single-slide native PPTX at out_path."""
    svg_path = Path(svg_path)
    svg_text = svg_path.read_text(encoding="utf-8")
    w, h = _viewbox_size(svg_text)
    svg_text = _preprocess_svg(svg_text, w, h)
    width_emu = int(round(w * EMU_PER_PX))
    height_emu = int(round(h * EMU_PER_PX))

    from pptx import Presentation  # imported lazily so module import works without it

    tmp = Path(tempfile.mkdtemp())
    try:
        prs = Presentation()
        prs.slide_width = width_emu
        prs.slide_height = height_emu
        prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        base_pptx = tmp / "base.pptx"
        prs.save(str(base_pptx))

        extract_dir = tmp / "pptx_content"
        with zipfile.ZipFile(base_pptx, "r") as zf:
            zf.extractall(extract_dir)

        # Converter reads from a file path; write the preprocessed SVG out.
        pre_svg = tmp / "source.svg"
        pre_svg.write_text(svg_text, encoding="utf-8")

        slide_xml, _media, _rels, _anim = convert_svg_to_slide_shapes(
            pre_svg, slide_num=1, verbose=verbose
        )
        # No media/rels for this skill's SVGs (no <image>); just overwrite slide1.xml.
        (extract_dir / "ppt" / "slides" / "slide1.xml").write_text(slide_xml, encoding="utf-8")

        out_path = Path(out_path)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for fp in extract_dir.rglob("*"):
                if fp.is_file():
                    zf.write(fp, fp.relative_to(extract_dir))
        return out_path
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Export architecture-diagram SVG/HTML to a native editable PPTX"
    )
    parser.add_argument("input", help="Path to .svg or .html file")
    parser.add_argument("-o", "--output", default=None,
                        help="Output .pptx path (default: <input>.pptx)")
    parser.add_argument("-v", "--verbose", action="store_true")
    args = parser.parse_args(argv)

    try:
        import pptx  # noqa: F401
    except ImportError:
        print("ERROR: python-pptx is not installed. Run: pip install -r requirements.txt",
              file=sys.stderr)
        return 2

    in_path = Path(args.input)
    if not in_path.exists():
        print(f"ERROR: input not found: {in_path}", file=sys.stderr)
        return 1

    if in_path.suffix.lower() in (".html", ".htm"):
        svg_text = _extract_svg_from_html(in_path.read_text(encoding="utf-8"))
        tmp_dir = Path(tempfile.mkdtemp())
        svg_path = tmp_dir / "extracted.svg"
        svg_path.write_text(svg_text, encoding="utf-8")
    else:
        svg_path = in_path

    out_path = Path(args.output) if args.output else in_path.with_suffix(".pptx")
    try:
        convert_svg_to_pptx(svg_path, out_path, verbose=args.verbose)
    except Exception as exc:  # noqa: BLE001
        print(f"ERROR: conversion failed: {exc}", file=sys.stderr)
        return 1
    print(f"Wrote: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
